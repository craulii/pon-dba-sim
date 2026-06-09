#!/usr/bin/env python3
"""
Genera la presentación de avances en PPTX para TEL-341 — OmneTeam.
Uso:  python3 generar_presentacion.py
Salida: presentacion_avances.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ─── Paleta de colores ──────────────────────────────────────────────────────
AZUL_USM    = RGBColor(0,   58,  122)   # Azul UTFSM institucional
AZUL_CLARO  = RGBColor(31,  119, 180)   # IPACT (color consistente con gráficos)
ROJO        = RGBColor(214,  39,  40)   # QoS-DBA (color consistente con gráficos)
BLANCO      = RGBColor(255, 255, 255)
GRIS        = RGBColor(90,   90,  90)
GRIS_CLARO  = RGBColor(240, 240, 245)
GRIS_MEDIO  = RGBColor(200, 200, 210)
NARANJA     = RGBColor(255, 127,  14)   # mMTC
VERDE       = RGBColor(44,  160,  44)   # OK / aprobado
AMARILLO    = RGBColor(255, 193,   7)   # Advertencia

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FIGURAS    = os.path.join(SCRIPT_DIR, "figuras")
OUT_FILE   = os.path.join(SCRIPT_DIR, "presentacion_avances.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── Helpers ────────────────────────────────────────────────────────────────

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, left, top, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, w, h,
        size=18, bold=False, italic=False,
        color=RGBColor(30, 30, 30),
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def txt_multi(slide, lines, left, top, w, h,
              size=16, color=RGBColor(30,30,30), leading=1.2):
    """lines: list of (text, bold, color_override_or_None)"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, bold, col) in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = col if col else color
    return tb


def img(slide, filename, left, top, w, h=None):
    path = os.path.join(FIGURAS, filename)
    if not os.path.exists(path):
        return None
    if h:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w))


def header_bar(slide, title, subtitle=None):
    """Barra superior azul con título y subtítulo opcionales."""
    rect(slide, 0, 0, 13.33, 1.3, AZUL_USM)
    txt(slide, title, 0.3, 0.08, 12.7, 0.7,
        size=28, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.75, 12.7, 0.45,
            size=14, bold=False, color=RGBColor(190, 210, 240),
            align=PP_ALIGN.LEFT)


def footer(slide, page_num):
    rect(slide, 0, 7.15, 13.33, 0.35, AZUL_USM)
    txt(slide, "OmneTeam  |  TEL-341 Simulación de Redes  |  UTFSM  |  2026",
        0.3, 7.17, 9, 0.3, size=9, color=RGBColor(190, 210, 240))
    txt(slide, f"{page_num}", 12.8, 7.17, 0.5, 0.3,
        size=9, color=RGBColor(190, 210, 240), align=PP_ALIGN.RIGHT)


def section_label(slide, text, color=AZUL_CLARO):
    """Etiqueta de sección pequeña arriba-derecha."""
    txt(slide, text, 10.5, 1.35, 2.6, 0.35,
        size=10, bold=True, color=color, align=PP_ALIGN.RIGHT)


def pill(slide, text, left, top, w=2.5, h=0.45,
         fill=AZUL_CLARO, text_color=BLANCO, size=14, bold=True):
    """Caja redondeada (pill) para destacar info."""
    r = slide.shapes.add_shape(9,   # RoundedRectangle
        Inches(left), Inches(top), Inches(w), Inches(h))
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    r.line.fill.background()
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return r


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
# Fondo bicolor
rect(sl, 0, 0, 13.33, 4.5, AZUL_USM)
rect(sl, 0, 4.5, 13.33, 3.0, GRIS_CLARO)

# Línea decorativa
rect(sl, 0, 4.35, 13.33, 0.15, AZUL_CLARO)

# Título principal
txt(sl, "Avances del Proyecto", 0.8, 0.8, 11.5, 1.0,
    size=38, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
txt(sl, "Evaluación de Algoritmos DBA en Redes PON", 0.8, 1.7, 11.5, 0.7,
    size=26, bold=False, color=RGBColor(190, 210, 240), align=PP_ALIGN.CENTER)
txt(sl, "bajo Tráfico 5G Multi-Servicio", 0.8, 2.3, 11.5, 0.6,
    size=26, bold=False, color=RGBColor(190, 210, 240), align=PP_ALIGN.CENTER)

# Detalles
txt(sl, "TEL-341 Simulación de Redes  —  Presentación de Avances Nº 2",
    1.0, 3.2, 11.0, 0.5, size=14, color=RGBColor(160, 190, 230), align=PP_ALIGN.CENTER)

# Pills de info
pill(sl, "OmneTeam", 1.5, 4.8, 2.5, 0.5, fill=AZUL_USM)
pill(sl, "David · José · Matías", 4.3, 4.8, 3.2, 0.5, fill=AZUL_CLARO)
pill(sl, "UTFSM  —  2026", 7.8, 4.8, 2.5, 0.5, fill=AZUL_USM)

# Descripción breve
txt(sl, "Simulación event-driven en OMNeT++ 6  •  C++17  •  Tráfico 5G (eMBB / URLLC / mMTC)",
    1.0, 5.7, 11.0, 0.5, size=13, color=GRIS, align=PP_ALIGN.CENTER)

footer(sl, "1")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OBJETIVO Y PROBLEMA
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Objetivo y Motivación", "¿Puede un algoritmo de scheduling proteger URLLC en un canal compartido?")
section_label(sl, "Contexto")

# Columna izquierda — problema
rect(sl, 0.3, 1.45, 5.8, 5.4, BLANCO,
     line_color=GRIS_MEDIO, line_width=0.5)
txt(sl, "El Problema", 0.5, 1.55, 5.4, 0.5, size=16, bold=True, color=AZUL_USM)
txt_multi(sl, [
    ("En 5G, las redes PON actúan como backhaul:", False, None),
    ("", False, None),
    ("  eMBB  →  streaming, alta tasa", False, AZUL_CLARO),
    ("  URLLC  →  latencia < 10 ms (crítico)", True, ROJO),
    ("  mMTC  →  millones de sensores IoT", False, NARANJA),
    ("", False, None),
    ("El canal upstream de una red PON es TDM:", False, None),
    ("solo transmite una ONU a la vez.", False, None),
    ("", False, None),
    ("IPACT (estándar) trata todos los paquetes", False, None),
    ("igual — no distingue eMBB de URLLC.", False, None),
    ("", False, None),
    ("Consecuencia: las ráfagas eMBB (Pareto)", False, None),
    ("agotan el canal y URLLC se descarta.", False, ROJO),
], 0.5, 2.1, 5.4, 4.5, size=14)

# Columna derecha — tabla 5G
rect(sl, 6.5, 1.45, 6.5, 5.4, BLANCO,
     line_color=GRIS_MEDIO, line_width=0.5)
txt(sl, "Clases de Tráfico 5G", 6.7, 1.55, 6.1, 0.5, size=16, bold=True, color=AZUL_USM)

rows = [
    ("Clase", "Distribución", "Requisito", AZUL_USM, BLANCO),
    ("eMBB",  "Pareto (α=1.5)", "Throughput alto", AZUL_CLARO, BLANCO),
    ("URLLC", "Poisson",        "Latencia < 10 ms", ROJO, BLANCO),
    ("mMTC",  "Periódico ±20%", "Baja tasa, escala", NARANJA, BLANCO),
]
row_h, col_w = 0.55, [1.5, 2.1, 2.1]
for i, (c1, c2, c3, bg, fg) in enumerate(rows):
    y = 2.2 + i * row_h
    for j, (text, w) in enumerate(zip([c1,c2,c3], col_w)):
        rect(sl, 6.7 + sum(col_w[:j]), y, w - 0.04, row_h - 0.04, bg)
        txt(sl, text, 6.75 + sum(col_w[:j]), y + 0.1, w - 0.1, row_h - 0.1,
            size=13, bold=(i==0), color=fg)

txt(sl, "★  Hipótesis: bajo carga real, IPACT falla en URLLC.",
    6.7, 4.55, 6.1, 0.4, size=13, bold=True, color=ROJO)
txt(sl, "★  QoS-DBA con prioridad estricta puede resolverlo.",
    6.7, 4.95, 6.1, 0.4, size=13, bold=True, color=VERDE)

footer(sl, "2")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARQUITECTURA
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Arquitectura de la Red Simulada",
           "OLT + Splitter 1:16 + 16 ONUs  |  Canal upstream TDM 1 Gbps")
section_label(sl, "Modelo")

# Diagrama de red (ASCII art visual con cajas)
# OLT
rect(sl, 0.4, 1.7, 2.2, 1.2, AZUL_USM)
txt(sl, "OLT", 0.5, 1.85, 2.0, 0.5, size=18, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
txt(sl, "Motor DBA", 0.5, 2.25, 2.0, 0.35, size=12, color=RGBColor(180,210,240), align=PP_ALIGN.CENTER)

# Línea fibra troncal
rect(sl, 2.6, 2.25, 1.5, 0.06, AZUL_CLARO)
txt(sl, "Fibra troncal\n100 μs", 2.6, 2.33, 1.5, 0.5, size=10, color=GRIS, align=PP_ALIGN.CENTER)

# Splitter
rect(sl, 4.1, 1.7, 1.7, 1.2, AZUL_CLARO)
txt(sl, "Splitter\n1:16", 4.1, 1.9, 1.7, 0.8, size=16, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# Líneas distribución + ONUs (simplificado — 4 representativas)
for i, label in enumerate(["ONU 0", "ONU 1", "...", "ONU 15"]):
    y = 1.5 + i * 1.1
    rect(sl, 5.8, y + 0.45, 1.2, 0.06, GRIS_MEDIO)
    rect(sl, 7.0, y + 0.15, 1.5, 0.8, AZUL_USM if label != "..." else GRIS_MEDIO)
    txt(sl, label, 7.0, y + 0.3, 1.5, 0.4, size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

txt(sl, "Distribución\n10 μs", 5.75, 3.1, 1.3, 0.5, size=10, color=GRIS, align=PP_ALIGN.CENTER)

# Columna derecha — parámetros
rect(sl, 8.9, 1.45, 4.1, 5.7, BLANCO, line_color=GRIS_MEDIO, line_width=0.5)
txt(sl, "Parámetros clave", 9.1, 1.55, 3.7, 0.45, size=15, bold=True, color=AZUL_USM)

params = [
    ("numONUs", "16"),
    ("dataRate", "1 Gbps"),
    ("pollingCycle", "2 ms"),
    ("maxGrantSize", "15 000 B"),
    ("guardTime", "1 μs"),
    ("urllcDeadline", "10 ms"),
    ("urllcRate", "5 Mbps / ONU"),
    ("mmtcRate", "0.5 Mbps / ONU"),
    ("embbRate", "50–200 Mbps / ONU"),
    ("Repeticiones", "3 (→ 10 pendiente)"),
]
for i, (k, v) in enumerate(params):
    y = 2.1 + i * 0.48
    txt(sl, k, 9.1, y, 2.0, 0.42, size=12, color=GRIS)
    txt(sl, v, 11.1, y, 1.8, 0.42, size=12, bold=True, color=AZUL_USM)

footer(sl, "3")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ALGORITMOS DBA
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Algoritmos DBA Implementados",
           "IPACT (estándar) vs QoS-DBA (propuesto)  —  intercambiables por parámetro")
section_label(sl, "DBA")

# Columna IPACT
rect(sl, 0.3, 1.45, 5.9, 5.7, BLANCO, line_color=AZUL_CLARO, line_width=1.5)
rect(sl, 0.3, 1.45, 5.9, 0.55, AZUL_CLARO)
txt(sl, "IPACT  —  Estándar IEEE 802.3ah", 0.5, 1.5, 5.5, 0.45,
    size=15, bold=True, color=BLANCO)

txt_multi(sl, [
    ("Polling proporcional sin QoS", True, AZUL_CLARO),
    ("", False, None),
    ("1. OLT envía POLL a todas las ONUs", False, None),
    ("2. ONU responde REPORT (total de bytes)", False, None),
    ("3. OLT asigna grant proporcional al total:", False, None),
    ("   grant_eMBB  = total × (eMBB_bytes/total)", False, GRIS),
    ("   grant_URLLC = total × (URLLC_bytes/total)", False, GRIS),
    ("", False, None),
    ("Problema:", True, ROJO),
    ("  Ráfaga Pareto eMBB → eMBB domina total", False, ROJO),
    ("  → URLLC recibe ~2 % del grant", False, ROJO),
    ("  → Cola URLLC crece, supera deadline", False, ROJO),
    ("  → Descarte masivo de URLLC", True, ROJO),
], 0.5, 2.1, 5.5, 4.7, size=13)

# Columna QoS-DBA
rect(sl, 6.9, 1.45, 6.1, 5.7, BLANCO, line_color=ROJO, line_width=1.5)
rect(sl, 6.9, 1.45, 6.1, 0.55, ROJO)
txt(sl, "QoS-DBA  —  Algoritmo Propuesto", 7.1, 1.5, 5.7, 0.45,
    size=15, bold=True, color=BLANCO)

txt_multi(sl, [
    ("Prioridad estricta + WFQ", True, ROJO),
    ("", False, None),
    ("Paso 1 — URLLC con prioridad absoluta:", True, ROJO),
    ("  grant_URLLC = min(URLLC_queue, cap)", False, None),
    ("  remaining   = cap − grant_URLLC", False, None),
    ("", False, None),
    ("Paso 2 — WFQ entre eMBB y mMTC:", True, AZUL_CLARO),
    ("  grant_eMBB  = remaining × 70 %", False, None),
    ("  grant_mMTC  = remaining × 30 %", False, None),
    ("  (ceder sobrante si hay poca demanda)", False, GRIS),
    ("", False, None),
    ("Garantía:", True, VERDE),
    ("  URLLC siempre transmite antes del", False, VERDE),
    ("  deadline, sin importar carga eMBB.", True, VERDE),
], 7.1, 2.1, 5.7, 4.7, size=13)

footer(sl, "4")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — IMPLEMENTACIÓN
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Módulos Implementados",
           "OMNeT++ 6.0  —  C++17  —  ~15 archivos  —  sin frameworks externos (INET)")
section_label(sl, "Código")

modules = [
    ("OLT.h / .cc",         "Motor central: envía POLLs, recibe REPORTs,\nejec. DBA, envía GRANTs. Registra cycleTime.", AZUL_USM),
    ("ONU.h / .cc",         "3 colas (eMBB/URLLC/mMTC), generadores de\ntráfico, registro de latencia/jitter/throughput.", AZUL_USM),
    ("IPACT.h / .cc",       "Algoritmo estándar: grant proporcional\nsin diferenciación de QoS.", AZUL_CLARO),
    ("QosDBA.h / .cc",      "Algoritmo propuesto: prioridad URLLC\n+ WFQ(70/30) para eMBB/mMTC.", ROJO),
    ("Splitter.h / .cc",    "Reenvío pasivo: recibe por un puerto,\ntransmite a todos los demás.", GRIS),
    ("PONMessages.msg",     "DataPacket · ReportMessage · GrantMessage\n(compilados auto. por opp_msgc a C++)", AZUL_CLARO),
    ("PONNetwork.ned",      "Topología parametrizable: OLT + splitter\n+ onu[numONUs]. Display strings para Qtenv.", AZUL_USM),
]

cols = 2
per_col = (len(modules) + 1) // cols
for i, (name, desc, color) in enumerate(modules):
    col = i // per_col
    row = i % per_col
    x = 0.3 + col * 6.55
    y = 1.5 + row * 0.83
    rect(sl, x, y, 6.2, 0.75, BLANCO, line_color=color, line_width=1)
    rect(sl, x, y, 0.25, 0.75, color)
    txt(sl, name, x + 0.35, y + 0.04, 2.5, 0.35, size=13, bold=True, color=color)
    txt(sl, desc, x + 0.35, y + 0.35, 5.7, 0.38, size=11, color=GRIS)

# Nota animación
rect(sl, 0.3, 6.85, 12.7, 0.45, AZUL_USM)
txt(sl, "★  GUI Qtenv: mensajes animados con colores (eMBB=verde, URLLC=rojo, mMTC=naranja, REPORT=azul, GRANT=cyan)  •  Bubbles en descarte",
    0.5, 6.88, 12.3, 0.38, size=11, bold=False, color=BLANCO)

footer(sl, "5")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ESTADO DE SIMULACIONES
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Estado de las Simulaciones",
           "24 corridas completadas  •  4 escenarios planificados  •  16 ONUs operativo")
section_label(sl, "Experimentos")

# Tabla de escenarios
headers = ["Configuración", "Algoritmo", "ONUs", "Carga eMBB", "Reps.", "Estado"]
col_w   = [2.5, 1.5, 0.8, 2.2, 0.8, 2.0]
rows_data = [
    ("IPACT_16ONU",  "IPACT",   "16", "50/100/150/200 Mbps", "3", "✓ Completado", VERDE),
    ("QoSDBA_16ONU", "QoS-DBA", "16", "50/100/150/200 Mbps", "3", "✓ Completado", VERDE),
    ("IPACT_32ONU",  "IPACT",   "32", "50/100/150/200 Mbps", "—", "Pendiente",     AMARILLO),
    ("QoSDBA_32ONU", "QoS-DBA", "32", "50/100/150/200 Mbps", "—", "Pendiente",     AMARILLO),
]

y_start = 1.55
# Cabecera
for j, (h, w) in enumerate(zip(headers, col_w)):
    rect(sl, 0.3 + sum(col_w[:j]), y_start, w - 0.05, 0.5, AZUL_USM)
    txt(sl, h, 0.37 + sum(col_w[:j]), y_start + 0.08, w - 0.1, 0.35,
        size=13, bold=True, color=BLANCO)

for i, (c1, c2, c3, c4, c5, c6, estado_color) in enumerate(rows_data):
    y = y_start + 0.5 + i * 0.6
    bg = BLANCO if i % 2 == 0 else GRIS_CLARO
    for j, (text, w) in enumerate(zip([c1,c2,c3,c4,c5,c6], col_w)):
        col_fill = estado_color if j == 5 else bg
        text_col = BLANCO if j == 5 else GRIS
        rect(sl, 0.3 + sum(col_w[:j]), y, w - 0.05, 0.55, col_fill,
             line_color=GRIS_MEDIO, line_width=0.3)
        txt(sl, text, 0.37 + sum(col_w[:j]), y + 0.1, w - 0.1, 0.38,
            size=12, bold=(j==5), color=text_col)

# Resumen de trabajo pendiente
rect(sl, 0.3, 5.1, 12.7, 1.9, BLANCO, line_color=AMARILLO, line_width=1)
txt(sl, "Trabajo pendiente para completar el proyecto:", 0.5, 5.2, 12.0, 0.4,
    size=14, bold=True, color=AZUL_USM)
txt_multi(sl, [
    ("①  Aumentar a 10 repeticiones → intervalos de confianza 95 % estadísticamente válidos", False, None),
    ("②  Ejecutar escenarios 32 ONUs → evaluar escalabilidad del comportamiento", False, None),
    ("③  Análisis estadístico completo: IC95 %, tests de hipótesis, verificación distribuciones", False, None),
    ("④  Informe final IEEE + demo GUI Qtenv", False, None),
], 0.5, 5.6, 12.3, 1.3, size=13)

footer(sl, "6")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RESULTADOS TABLA
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Resultados Preliminares — Tasa de Pérdida",
           "Métrica más importante: lossURLLC  |  Carga 50–200 Mbps  |  16 ONUs  |  3 repeticiones")
section_label(sl, "Resultados")

# Tabla de resultados
headers2 = ["Algoritmo", "Carga eMBB", "Pérdida URLLC", "Pérdida eMBB", "Pérdida mMTC"]
col_w2   = [2.0, 1.8, 2.5, 2.0, 2.0]
ipact_rows = [
    ("IPACT",   "50 Mbps",  "75.8 %",  "0.1 %",  "≈ 0 %"),
    ("IPACT",   "100 Mbps", "86.8 %",  "49.0 %", "≈ 0 %"),
    ("IPACT",   "150 Mbps", "86.9 %",  "66.0 %", "≈ 0 %"),
    ("IPACT",   "200 Mbps", "87.0 %",  "74.5 %", "≈ 0 %"),
    ("QoS-DBA", "50 Mbps",  "0.0 %",   "8.1 %",  "≈ 0 %"),
    ("QoS-DBA", "100 Mbps", "0.0 %",   "53.9 %", "≈ 0 %"),
    ("QoS-DBA", "150 Mbps", "0.0 %",   "69.2 %", "≈ 0 %"),
    ("QoS-DBA", "200 Mbps", "0.0 %",   "76.9 %", "≈ 0 %"),
]

y_s = 1.55
for j, (h, w) in enumerate(zip(headers2, col_w2)):
    rect(sl, 0.3 + sum(col_w2[:j]), y_s, w - 0.05, 0.5, AZUL_USM)
    txt(sl, h, 0.37 + sum(col_w2[:j]), y_s + 0.08, w - 0.1, 0.35,
        size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

for i, row in enumerate(ipact_rows):
    y = y_s + 0.5 + i * 0.52
    is_ipact = row[0] == "IPACT"
    bg = RGBColor(235, 242, 252) if is_ipact else RGBColor(252, 237, 237)
    for j, (text, w) in enumerate(zip(row, col_w2)):
        if j == 2:  # lossURLLC — resaltar
            cell_fill = ROJO if is_ipact else VERDE
            text_col  = BLANCO
            bold      = True
        else:
            cell_fill = bg
            text_col  = GRIS
            bold      = False
        rect(sl, 0.3 + sum(col_w2[:j]), y, w - 0.05, 0.5, cell_fill,
             line_color=GRIS_MEDIO, line_width=0.3)
        txt(sl, text, 0.37 + sum(col_w2[:j]), y + 0.1, w - 0.1, 0.33,
            size=12, bold=bold, color=text_col, align=PP_ALIGN.CENTER)

# Insight clave
rect(sl, 0.3, 6.65, 12.7, 0.55, AZUL_USM)
txt(sl, "★  Con IPACT: 3 de cada 4 paquetes URLLC se pierden incluso a carga baja  "
        "|  Con QoS-DBA: pérdida URLLC = 0 % en todos los escenarios",
    0.5, 6.7, 12.3, 0.45, size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

footer(sl, "7")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GRÁFICO: PÉRDIDA POR CLASE (explicado)
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Gráfico: Tasa de Pérdida por Clase",
           "Escala logarítmica  |  Carga eMBB = 200 Mbps  |  Diferencia de 5 órdenes de magnitud en URLLC")
section_label(sl, "Gráfico 1/4")

img(sl, "packet_loss_by_class.png", 0.3, 1.45, 7.5, 5.5)

# Anotaciones explicativas
rect(sl, 8.2, 1.45, 4.8, 5.5, BLANCO, line_color=GRIS_MEDIO, line_width=0.5)
txt(sl, "Cómo leer este gráfico", 8.4, 1.55, 4.4, 0.45, size=15, bold=True, color=AZUL_USM)
txt_multi(sl, [
    ("Eje Y en escala logarítmica:", True, AZUL_USM),
    ("  Permite ver diferencias de muchos", False, None),
    ("  órdenes de magnitud en una figura.", False, None),
    ("", False, None),
    ("eMBB (izquierda):", True, AZUL_CLARO),
    ("  Ambos ~75–77 %. Canal saturado.", False, None),
    ("  La diferencia entre algoritmos es mínima.", False, None),
    ("", False, None),
    ("URLLC (centro):", True, ROJO),
    ("  IPACT → ~87 % (barra azul llena)", False, ROJO),
    ("  QoS-DBA → ~0 % (barra roja al piso)", False, VERDE),
    ("  Diferencia: 5 órdenes de magnitud.", True, ROJO),
    ("", False, None),
    ("mMTC (derecha):", True, NARANJA),
    ("  Ambos ~0 %. Tasa baja, sin problema.", False, None),
    ("", False, None),
    ("Conclusión:", True, VERDE),
    ("  QoS-DBA elimina completamente", True, VERDE),
    ("  la pérdida de tráfico crítico.", True, VERDE),
], 8.4, 2.1, 4.4, 4.6, size=12)

footer(sl, "8")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — GRÁFICO: CDF LATENCIA URLLC (explicado)
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Gráfico: CDF de Latencia URLLC",
           "Distribución acumulada  |  Carga 200 Mbps  |  Eje X logarítmico")
section_label(sl, "Gráfico 2/4")

img(sl, "cdf_latency_urllc.png", 0.3, 1.45, 7.5, 5.5)

rect(sl, 8.2, 1.45, 4.8, 5.5, BLANCO, line_color=GRIS_MEDIO, line_width=0.5)
txt(sl, "Cómo leer este gráfico", 8.4, 1.55, 4.4, 0.45, size=15, bold=True, color=AZUL_USM)
txt_multi(sl, [
    ("La CDF muestra qué fracción de", False, None),
    ("paquetes tiene latencia ≤ X μs.", False, None),
    ("", False, None),
    ("QoS-DBA (rojo):", True, ROJO),
    ("  Curva suave 100–3 000 μs.", False, None),
    ("  P50 ≈ 1 500 μs (1.5 ms)", False, None),
    ("  P99 ≈ 2 500 μs (2.5 ms)", False, None),
    ("  100% dentro del deadline (10 ms).", True, VERDE),
    ("", False, None),
    ("IPACT (azul):", True, AZUL_CLARO),
    ("  Salto brusco en ~9 000–10 000 μs.", False, None),
    ("  Solo los paquetes que 'sobreviven'", False, None),
    ("  lo hacen al límite del deadline.", False, ROJO),
    ("  El 87% fue descartado antes.", True, ROJO),
    ("", False, None),
    ("La línea vertical roja es el deadline", False, None),
    ("(10 ms = 10 000 μs).", False, None),
    ("QoS-DBA opera 4× por debajo.", True, VERDE),
], 8.4, 2.1, 4.4, 4.6, size=12)

footer(sl, "9")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — GRÁFICO: LATENCIA PROMEDIO Y P99 (explicado)
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Latencia Promedio y P99 URLLC",
           "eMBB/URLLC/mMTC vs carga  |  QoS-DBA reduce latencia URLLC en 10×")
section_label(sl, "Gráfico 3/4")

img(sl, "latency_avg_by_class.png",    0.2, 1.4, 6.2, 4.6)
img(sl, "latency_p99_urllc_vs_load.png", 6.5, 1.4, 6.6, 4.6)

# Notas debajo
rect(sl, 0.2, 6.0, 6.2, 1.2, BLANCO, line_color=GRIS_MEDIO, line_width=0.5)
txt_multi(sl, [
    ("Latencia promedio (izquierda):", True, AZUL_USM),
    ("URLLC: IPACT ~10 ms  vs  QoS-DBA ~1 ms  (10× mejor)", False, None),
    ("mMTC: IPACT ~85 ms  vs  QoS-DBA ~1 ms  (85× mejor)", False, None),
    ("eMBB mayor en QoS-DBA: sesgo de supervivencia (más paquetes sirven, esperan más)", False, GRIS),
], 0.4, 6.1, 5.8, 1.05, size=11)

rect(sl, 6.5, 6.0, 6.6, 1.2, BLANCO, line_color=GRIS_MEDIO, line_width=0.5)
txt_multi(sl, [
    ("P99 URLLC vs carga (derecha):", True, AZUL_USM),
    ("IPACT: P99 = 10 ms en todos los niveles → siempre al límite del deadline", False, ROJO),
    ("QoS-DBA: P99 ≈ 2.4 ms constante → independiente de la carga eMBB", False, VERDE),
    ("Ambas curvas son planas: el mecanismo DBA domina, no la carga ofrecida", False, GRIS),
], 6.7, 6.1, 6.2, 1.05, size=11)

footer(sl, "10")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DASHBOARD (slide para presentación)
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Dashboard — Comparación Completa",
           "4 métricas en una figura  |  Resultado central del proyecto")
section_label(sl, "Resumen")

img(sl, "summary_dashboard.png", 0.3, 1.35, 12.7, 5.9)

footer(sl, "11")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PRÓXIMOS PASOS
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(sl, "Próximos Pasos",
           "Para completar el proyecto y el informe final")
section_label(sl, "Pendiente")

steps = [
    ("1", "10 repeticiones",
     "Ampliar de 3 a 10 corridas por escenario.\nPermite IC 95 % estadísticamente válidos.\nEstimado: ~15 min de cómputo adicional.",
     AZUL_USM, "Alta"),
    ("2", "Escenarios 32 ONUs",
     "Ejecutar IPACT_32ONU y QoSDBA_32ONU.\nEvalúa escalabilidad: ¿se mantiene la\nprotección URLLC con más usuarios?",
     AZUL_CLARO, "Alta"),
    ("3", "Análisis estadístico",
     "IC 95 % con t-Student, tests Mann-Whitney U,\nverificación de distribuciones (KS test),\nresultados con significancia estadística.",
     ROJO, "Media"),
    ("4", "Informe + Demo GUI",
     "Informe final estilo IEEE con todas\nlas secciones formales. Demo Qtenv\ncon animación en vivo para presentación.",
     VERDE, "Baja"),
]

for i, (num, title, desc, color, prioridad) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.55
    y = 1.5 + row * 2.6

    rect(sl, x, y, 6.2, 2.4, BLANCO, line_color=color, line_width=1.5)
    rect(sl, x, y, 0.7, 2.4, color)

    txt(sl, num, x + 0.05, y + 0.85, 0.6, 0.7, size=22, bold=True,
        color=BLANCO, align=PP_ALIGN.CENTER)
    txt(sl, title, x + 0.85, y + 0.08, 5.1, 0.55, size=16, bold=True, color=color)

    pill(sl, f"Prioridad {prioridad}", x + 0.85, y + 0.62, 2.0, 0.35,
         fill=color, size=10)

    txt(sl, desc, x + 0.85, y + 1.05, 5.1, 1.25, size=12, color=GRIS)

footer(sl, "12")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CIERRE
# ═══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
rect(sl, 0, 0, 13.33, 7.5, AZUL_USM)
rect(sl, 0, 4.8, 13.33, 2.7, GRIS_CLARO)
rect(sl, 0, 4.65, 13.33, 0.15, AZUL_CLARO)

txt(sl, "Conclusión Preliminar", 1.0, 0.9, 11.2, 0.9,
    size=32, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

txt_multi(sl, [
    ("✓   IPACT falla en proteger URLLC: 75–87 % de pérdida desde carga moderada", True, BLANCO),
    ("✓   QoS-DBA garantiza pérdida URLLC = 0 % en todos los niveles de carga",   True, RGBColor(144,238,144)),
    ("✓   El trade-off: eMBB pierde ~5 pp adicionales — el precio esperado por proteger URLLC", False, RGBColor(190,210,240)),
    ("✓   El modelo compila, corre y produce resultados coherentes con la teoría",  False, RGBColor(190,210,240)),
], 1.0, 1.9, 11.2, 2.5, size=16)

txt(sl, "¿Preguntas?", 1.0, 5.0, 11.2, 1.0,
    size=36, bold=True, color=AZUL_USM, align=PP_ALIGN.CENTER)

txt(sl, "OmneTeam  ·  David Retuerto  ·  José Vega  ·  Matías Perelli",
    1.0, 6.2, 11.2, 0.5, size=14, color=GRIS, align=PP_ALIGN.CENTER)
txt(sl, "TEL-341 Simulación de Redes  ·  UTFSM  ·  2026",
    1.0, 6.6, 11.2, 0.4, size=12, color=GRIS, align=PP_ALIGN.CENTER)

# No footer en slide de cierre


# ─── Guardar ────────────────────────────────────────────────────────────────
prs.save(OUT_FILE)
print(f"Presentación guardada en: {OUT_FILE}")
print(f"Total de slides: {len(prs.slides)}")
